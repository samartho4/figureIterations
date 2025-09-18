#!/usr/bin/env python3
"""
make_presentation.py — builds a PPTX deck from clean_figures_final + paper numbers.

Usage:
  python scripts/make_presentation.py --fig-dir clean_figures_final --out slides/microgrid_ude_bnode_neurips.pptx
"""

import argparse, glob, os, sys, datetime, subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TITLE = "Learning Microgrid Dynamics via UDEs and Bayesian Neural ODEs"
SUBTITLE = "Accurate, interpretable, and calibrated models for microgrid decision-making"

# Paper numbers (kept here so they stay consistent everywhere)
NUM = {
    "test_scenarios": 10,
    "test_points": 2010,
    "ude_rmse_x2": 0.2475,
    "phys_rmse_x2": 0.2520,
    "wilcoxon_p": 0.9219,
    "delta_mean": -0.004488,
    "ci_low": -0.038517,
    "ci_high": 0.031438,
    "cov50": 0.541,
    "cov90": 0.849,
    "nll_raw": 268800.794,
    "nll_cal": 4088.593,
    "nll_red_pct": 98.48,
    "r2_symbolic": 0.982,
    "cubic_c": 0.018945,
}

# Figures to place (search by basename; we’ll pick .png or .jpg)
# Adjust the "title" or "subtitle" if you want different slide headings.
FIGS = [
    {"basename": "reliability_diagram",       "title": "BNODE Reliability", "subtitle": "Empirical vs nominal coverage (square plot)"},
    {"basename": "posterior_predictive_checks","title": "Posterior Predictive Checks", "subtitle": "Trajectories with median, 50/90% bands"},
    {"basename": "calibration_error",         "title": "Calibration Sweep", "subtitle": "Coverage error & NLL vs α_cal"},
    {"basename": "ude_residual_cubic",        "title": "UDE Residual (Symbolic)", "subtitle": "Cubic fit; R² ≈ %.3f" % NUM["r2_symbolic"]},
    {"basename": "ude_ablations",             "title": "UDE Ablations", "subtitle": "Width & λ sweeps"},
    {"basename": "noise_robustness",          "title": "Noise Robustness", "subtitle": "UDE RMSE & BNODE intervals vs σ"},
]

# Speaker notes per slide (index-aligned with SLIDES spec below)
NOTES = {
"motivation": """Microgrids face fast, nonlinear, mode-switching dynamics. High-fidelity DAEs are stiff and slow; black-box models break physics and fail OOD. We use SciML: UDE keeps physics and learns the residual; BNODE learns full dynamics with calibrated uncertainty.""",

"limits": """Why not PINNs? Known failure modes (trivial minima, gradient pathologies, spectral bias) and trouble with stiffness. UDE/BNODE preserve structure and use proper ODE solvers and Bayesian inference.""",

"model": """Two-state model: x1 = storage SoC with efficiency; x2 = frequency/power deviation with damping, power coupling, and storage interaction. This minimal model captures key operational behavior.""",

"ude": """UDE learns only f_theta(P_gen), preserving storage physics. Tiny network (n=3; 9 params), Rosenbrock23 for stiffness, composite loss focused on x2. Symbolic extraction later shows a cubic saturation.""",

"bnode": """BNODE places distributions over ODE parameters. We run NUTS, check R-hat/ESS, then apply a single α_cal variance scale to hit nominal coverage. Gives risk-aware predictions useful for chance constraints.""",

"metrics": """Test set: %d scenarios, %d points. UDE vs physics on x2 RMSE: %.4f vs %.4f; Δ≈%.6f (95%% CI [%.6f, %.6f]), Wilcoxon p=%.4f. BNODE calibration: 50%%/90%% ≈ %.3f/%.3f; NLL drops ~%.2f%% after calibration.""" % (
    NUM["test_scenarios"], NUM["test_points"],
    NUM["ude_rmse_x2"], NUM["phys_rmse_x2"],
    NUM["delta_mean"], NUM["ci_low"], NUM["ci_high"], NUM["wilcoxon_p"],
    NUM["cov50"], NUM["cov90"], NUM["nll_red_pct"]
),

"reliability": """Post-calibration, reliability curves hug the diagonal. Slight over/under-confidence removed by a single global variance scale α_cal. This is actionable for risk-aware ops.""",

"ppc": """Posterior predictive bands cover trajectories without bias. No post-warmup divergences, good ESS — posterior seems healthy.""",

"calib": """Global α_cal ≈ 1.8 gives minimal coverage error and massive NLL improvement (%.0f → %.0f). One knob; big payoff.""" % (NUM["nll_raw"], NUM["nll_cal"]),

"symbolic": """UDE residual cubic fit (R²≈%.3f): mild saturation as P_gen→1. Suggests adaptive droop: β(P_gen)=β0(1 - c·P_gen²), with c≈%.6f. Physically interpretable and useful for controller tuning."""
            % (NUM["r2_symbolic"], NUM["cubic_c"]),

"ablations": """UDE is stable across width and λ; width=3 suffices. Regularization trades fit vs generalization smoothly.""",

"robustness": """As input noise rises, UDE RMSE degrades ~linearly; BNODE widens intervals to maintain coverage — distinguishes epistemic vs aleatoric.""",

"runtime": """UDE stays fast (sub-ms per trajectory) with a stiff solver in the loop, suitable for real-time or MPC contexts.""",

"takeaway": """Use UDE for real-time, physics-faithful control; use BNODE for planning under uncertainty. Together: interpretable, calibrated, and operationally ready."""
}

def find_image(fig_dir, base):
    for ext in (".png", ".jpg", ".jpeg"):
        path = os.path.join(fig_dir, f"{base}_improved{ext}")
        if os.path.exists(path):
            return path
        # also accept files like <base>.png if *_improved not found
        path2 = os.path.join(fig_dir, f"{base}{ext}")
        if os.path.exists(path2):
            return path2
    return None

def add_footer(slide, text):
    txbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.3))
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(10)
    font.color.rgb = RGBColor(120,120,120)

def add_title_and_bullets(prs, title, bullets, notes=None, footer=None):
    slide_layout = prs.slide_layouts[1]  # Title + Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i>0 else body.paragraphs[0]
        p.text = b
        p.level = 0
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    if footer:
        add_footer(slide, footer)
    return slide

def add_image_slide(prs, title, subtitle, img_path, notes=None, footer=None):
    layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title if subtitle is None else f"{title} — {subtitle}"
    if img_path and os.path.exists(img_path):
        pic = slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.5))
        # scale to fit within (width=12.5in, height=5.5in area)
        max_w, max_h = Inches(12.5-1.2), Inches(5.5)
        scale = min(max_w/pic.width, max_h/pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        # center horizontally
        pic.left = int((prs.slide_width - pic.width) / 2)
    else:
        tf = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(10), Inches(1.0)).text_frame
        tf.text = f"[Missing image: {img_path}]"
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    if footer:
        add_footer(slide, footer)
    return slide

def git_short_sha():
    try:
        out = subprocess.check_output(["git", "rev-parse", "--short", "HEAD"], stderr=subprocess.DEVNULL)
        return out.decode().strip()
    except Exception:
        return "unversioned"

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--fig-dir", required=True, help="Path to clean_figures_final")
    ap.add_argument("--out", required=True, help="Output PPTX path")
    args = ap.parse_args()

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)  # 16:9 widescreen

    # Footer provenance
    now = datetime.datetime.now().strftime("%Y-%m-%d")
    sha = git_short_sha()
    footer = f"Provenance: release/minimal-figures-and-analysis @ {sha} • generated {now}"

    # --- Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title
    slide.shapes.title.text = TITLE
    slide.placeholders[1].text = SUBTITLE
    add_footer(slide, footer)

    # --- Motivation & problem
    add_title_and_bullets(
        prs,
        "Why SciML for Microgrids?",
        [
            "Fast, nonlinear, mode-switching dynamics; stiff DAEs are slow, black-box learners violate physics",
            "Goal: accurate, interpretable, calibrated surrogates for control & what-if planning",
            "Approach: UDE (physics + learned residual) and BNODE (Bayesian Neural ODE with calibration)",
        ],
        notes=NOTES["motivation"],
        footer=footer
    )

    # --- Limits of existing approaches
    add_title_and_bullets(
        prs,
        "Limits of PINNs & Pure Black-Box",
        [
            "PINNs: failure modes & gradient pathologies; spectral bias; stiffness challenges",
            "Black-box: poor extrapolation; constraint violations",
            "Structure-preserving ODEs + stiff solvers avoid these issues",
        ],
        notes=NOTES["limits"],
        footer=footer
    )

    # --- Two-state model
    add_title_and_bullets(
        prs,
        "Two-State Microgrid Model",
        [
            "x₁: storage SoC with efficiency losses;  x₂: frequency/power deviation",
            "dx₂/dt = −α x₂ + β(P_gen − P_load) + γ x₁",
            "Testbed for structure-aware learning and calibrated UQ",
        ],
        notes=NOTES["model"],
        footer=footer
    )

    # --- UDE slide
    add_title_and_bullets(
        prs,
        "UDE: Physics + Learned Residual",
        [
            "Learn fθ(P_gen) only; preserve storage physics",
            "Tiny MLP (n=3; 9 params), Rosenbrock23, composite loss favoring x₂",
            "Symbolic extraction → cubic saturation insight",
        ],
        notes=NOTES["ude"],
        footer=footer
    )

    # --- BNODE slide
    add_title_and_bullets(
        prs,
        "BNODE: Uncertainty-Aware Dynamics",
        [
            "Bayesian parameters; NUTS sampling; PPC checks",
            "Single α_cal post-calibration for nominal coverage",
            "Enables chance-constrained, risk-aware ops",
        ],
        notes=NOTES["bnode"],
        footer=footer
    )

    # --- Metrics summary
    add_title_and_bullets(
        prs,
        "Results Summary",
        [
            f"Test set: {NUM['test_scenarios']} scenarios / {NUM['test_points']} points",
            f"UDE vs Physics (x₂ RMSE): {NUM['ude_rmse_x2']:.4f} vs {NUM['phys_rmse_x2']:.4f}",
            f"Δ={NUM['delta_mean']:.6f} (95% CI [{NUM['ci_low']:.6f}, {NUM['ci_high']:.6f}]), Wilcoxon p={NUM['wilcoxon_p']:.4f}",
            f"BNODE coverage ~ {NUM['cov50']:.3f} / {NUM['cov90']:.3f} (50% / 90%), NLL ↓ ~{NUM['nll_red_pct']:.2f}% after calibration",
        ],
        notes=NOTES["metrics"],
        footer=footer
    )

    # --- Figure slides
    for spec in FIGS:
        img = find_image(args.fig_dir, spec["basename"])
        notes_key = {
            "reliability_diagram": "reliability",
            "posterior_predictive_checks": "ppc",
            "calibration_error": "calib",
            "ude_residual_cubic": "symbolic",
            "ude_ablations": "ablations",
            "noise_robustness": "robustness",
        }.get(spec["basename"], None)
        add_image_slide(
            prs,
            spec["title"],
            spec["subtitle"],
            img,
            notes=NOTES.get(notes_key),
            footer=footer
        )

    # --- Runtime / Deployment
    add_title_and_bullets(
        prs,
        "Runtime & Deployment",
        [
            "Sub-ms UDE trajectories with stiff solver → real-time feasible",
            "BNODE for planning: intervals for chance constraints",
            "Hybrid strategy: UDE in the loop, BNODE audits",
        ],
        notes=NOTES["runtime"],
        footer=footer
    )

    # --- Takeaways
    add_title_and_bullets(
        prs,
        "Takeaways",
        [
            "UDE: physics-comparable accuracy + interpretability (cubic saturation)",
            "BNODE: calibrated UQ (single α_cal), actionable risk",
            "Together: reliable surrogates for control & planning",
        ],
        notes=NOTES["takeaway"],
        footer=footer
    )

    os.makedirs(os.path.dirname(args.out), exist_ok=True)
    prs.save(args.out)
    print(f"Wrote {args.out}")

if __name__ == "__main__":
    main()
