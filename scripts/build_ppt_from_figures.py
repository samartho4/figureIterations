# scripts/build_ppt_from_figures.py
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from scripts.fig_manifest import FIGS

SRC = Path("clean_figures_final")
OUT = Path("slides")
OUT.mkdir(exist_ok=True)

def add_title(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_image_slide(prs, title, caption, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only
    slide.shapes.title.text = title
    left = Inches(0.6); top = Inches(1.2); width = Inches(11.2)
    slide.shapes.add_picture(str(img_path), left, top, width=width)
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(11.2), Inches(1.0))
    p = tx.text_frame.paragraphs[0]
    p.text = caption; p.font.size = Pt(14); p.alignment = PP_ALIGN.LEFT


def main():
    prs = Presentation()
    add_title(prs, "Learning Microgrid Dynamics via UDEs & BNODEs",
              "Accurate, interpretable, and calibrated models for microgrid decision-making")

    # Add slides for every *_improved.png we actually have
    for name, (title, caption) in FIGS.items():
        img = SRC / f"{name}.png"
        if img.exists():
            add_image_slide(prs, title, caption, img)
        else:
            print("[skip] missing", img)
    out = OUT/"microgrid_ude_bnode_neurips.pptx"
    prs.save(out)
    print("[ok] wrote", out)


if __name__ == "__main__":
    main()
